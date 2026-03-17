"""
update_pptx.py
Update Hospitalization_IO.pptx with new figures and table images.
"""
import sys
sys.path.insert(0, r'C:\users\hsaee\desktop\cms_viewer\env\Lib\site-packages')

import io
import openpyxl
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib import rcParams
from pptx import Presentation
from pptx.util import Inches, Emu

OUT_DIR  = r"C:\Users\hsaee\Desktop\CMS_viewer\projects\HNC_io_hosp"
PPTX_IN  = rf"{OUT_DIR}\Hospitalization_IO.pptx"
PPTX_OUT = rf"{OUT_DIR}\Hospitalization_IO.pptx"

PRIMARY  = '#BA0C2F'
DARK     = '#7A0820'
TINT     = '#F5D0D6'
ALT      = '#F9ECEE'
HEADER_BG = '#BA0C2F'
HEADER_FG = 'white'

# ── Helper: remove shape and add picture at same position/size ─────────────────
def replace_shape_with_image(slide, shape_name, img_path):
    for shape in slide.shapes:
        if shape.name == shape_name:
            left   = shape.left
            top    = shape.top
            width  = shape.width
            height = shape.height
            sp = shape._element
            sp.getparent().remove(sp)
            slide.shapes.add_picture(img_path, left, top, width, height)
            print(f"  Replaced '{shape_name}' with {img_path}")
            return True
    print(f"  WARNING: Shape '{shape_name}' not found.")
    return False

# ── Render table from xlsx as PNG ─────────────────────────────────────────────
def render_table1_png(xlsx_path, out_png, fig_width=5.8, fig_height=6.22):
    """Render Table 1 (patient characteristics) as a PNG image."""
    wb = openpyxl.load_workbook(xlsx_path)
    ws = wb.active

    # Collect all rows
    rows = []
    for row in ws.iter_rows(values_only=True):
        rows.append(row)

    # Find header row (first row with ≥4 non-empty cells)
    header_idx = None
    for i, row in enumerate(rows):
        filled = sum(1 for c in row if c is not None and str(c).strip())
        if filled >= 4:
            header_idx = i
            break
    if header_idx is None:
        header_idx = 2

    title_text = rows[0][0] if rows else ''
    headers = [str(c) if c is not None else '' for c in rows[header_idx][:5]]
    data_rows = []
    section_flags = []

    # Read each cell to detect section styling
    for row_idx in range(header_idx + 1, len(rows)):
        row = rows[row_idx]
        if all(c is None for c in row):
            continue
        vals = [str(c) if c is not None else '' for c in row[:5]]
        # Detect section rows: bold + SECTION_FILL fill
        cell = ws.cell(row=row_idx + 1, column=1)
        is_section = (cell.font and cell.font.bold and
                      cell.fill and cell.fill.fill_type == 'solid' and
                      cell.fill.fgColor.rgb in ('F5D0D6', '00F5D0D6'))
        data_rows.append(vals)
        section_flags.append(is_section)

    # Build figure
    rcParams['font.family'] = 'sans-serif'
    fig, ax = plt.subplots(figsize=(fig_width, fig_height))
    fig.patch.set_facecolor('white')
    ax.axis('off')

    # Title
    ax.text(0.5, 0.99, title_text, transform=ax.transAxes,
            ha='center', va='top', fontsize=7.5, fontweight='bold',
            color=PRIMARY, wrap=True)

    # Column widths (proportional)
    col_widths = [0.38, 0.16, 0.16, 0.16, 0.08]
    n_cols = 5

    # Draw header
    HEADER_H = 0.95
    TITLE_H  = 0.01
    ROW_H = (HEADER_H - 0.04) / max(len(data_rows) + 1, 1)

    x_starts = [sum(col_widths[:i]) for i in range(n_cols)]
    y_header = HEADER_H - ROW_H

    for ci, (hdr, x, w) in enumerate(zip(headers, x_starts, col_widths)):
        rect = mpatches.FancyBboxPatch(
            (x, y_header), w, ROW_H - 0.002,
            boxstyle='square,pad=0', facecolor=HEADER_BG,
            edgecolor='white', linewidth=0.3, transform=ax.transAxes, clip_on=False)
        ax.add_patch(rect)
        ax.text(x + w/2, y_header + ROW_H/2, hdr,
                transform=ax.transAxes,
                ha='center', va='center', fontsize=5.5, fontweight='bold',
                color='white', wrap=True)

    # Draw data rows
    for ri, (row, is_sec) in enumerate(zip(data_rows, section_flags)):
        y = y_header - (ri + 1) * ROW_H
        if y < 0:
            break
        bg = TINT if is_sec else (ALT if ri % 2 == 0 else 'white')
        for ci, (val, x, w) in enumerate(zip(row, x_starts, col_widths)):
            rect = mpatches.FancyBboxPatch(
                (x, y), w, ROW_H - 0.001,
                boxstyle='square,pad=0', facecolor=bg,
                edgecolor='#dddddd', linewidth=0.2,
                transform=ax.transAxes, clip_on=False)
            ax.add_patch(rect)
            fs = 5.0 if is_sec else 4.8
            fw = 'bold' if is_sec else 'normal'
            fc = DARK if is_sec else '#222222'
            ha = 'left' if ci == 0 else 'center'
            pad = 0.005 if ci == 0 else 0
            ax.text(x + pad, y + ROW_H/2, val[:60],
                    transform=ax.transAxes,
                    ha=ha, va='center', fontsize=fs, fontweight=fw,
                    color=fc, clip_on=True)

    plt.savefig(out_png, dpi=300, bbox_inches='tight', facecolor='white',
                edgecolor='none', pad_inches=0.02)
    plt.close()
    print(f"  Rendered table: {out_png}")


def render_table2_png(xlsx_path, out_png, fig_width=6.43, fig_height=4.0):
    """Render Table 2 (hospice utilization) as a PNG image."""
    wb = openpyxl.load_workbook(xlsx_path)
    ws = wb.active

    rows = []
    for row in ws.iter_rows(values_only=True):
        rows.append(row)

    header_idx = None
    for i, row in enumerate(rows):
        filled = sum(1 for c in row if c is not None and str(c).strip())
        if filled >= 2:
            header_idx = i
            break
    if header_idx is None:
        header_idx = 2

    title_text = rows[0][0] if rows else ''
    headers = [str(c) if c is not None else '' for c in rows[header_idx][:3]]
    data_rows = []
    section_flags = []

    for row_idx in range(header_idx + 1, len(rows)):
        row = rows[row_idx]
        if all(c is None for c in row):
            continue
        vals = [str(c) if c is not None else '' for c in row[:3]]
        cell = ws.cell(row=row_idx + 1, column=1)
        is_section = (cell.font and cell.font.bold and
                      cell.fill and cell.fill.fill_type == 'solid' and
                      cell.fill.fgColor.rgb in ('F5D0D6', '00F5D0D6'))
        data_rows.append(vals)
        section_flags.append(is_section)

    rcParams['font.family'] = 'sans-serif'
    fig, ax = plt.subplots(figsize=(fig_width, fig_height))
    fig.patch.set_facecolor('white')
    ax.axis('off')

    ax.text(0.5, 0.99, title_text, transform=ax.transAxes,
            ha='center', va='top', fontsize=7.5, fontweight='bold',
            color=PRIMARY, wrap=True)

    col_widths = [0.58, 0.28, 0.14]
    n_cols = 3

    HEADER_H = 0.93
    ROW_H = (HEADER_H - 0.04) / max(len(data_rows) + 1, 1)

    x_starts = [sum(col_widths[:i]) for i in range(n_cols)]
    y_header = HEADER_H - ROW_H

    for ci, (hdr, x, w) in enumerate(zip(headers, x_starts, col_widths)):
        rect = mpatches.FancyBboxPatch(
            (x, y_header), w, ROW_H - 0.002,
            boxstyle='square,pad=0', facecolor=HEADER_BG,
            edgecolor='white', linewidth=0.3, transform=ax.transAxes, clip_on=False)
        ax.add_patch(rect)
        ax.text(x + w/2, y_header + ROW_H/2, hdr,
                transform=ax.transAxes,
                ha='center', va='center', fontsize=5.5, fontweight='bold',
                color='white', wrap=True)

    for ri, (row, is_sec) in enumerate(zip(data_rows, section_flags)):
        y = y_header - (ri + 1) * ROW_H
        if y < 0:
            break
        bg = TINT if is_sec else (ALT if ri % 2 == 0 else 'white')
        for ci, (val, x, w) in enumerate(zip(row, x_starts, col_widths)):
            rect = mpatches.FancyBboxPatch(
                (x, y), w, ROW_H - 0.001,
                boxstyle='square,pad=0', facecolor=bg,
                edgecolor='#dddddd', linewidth=0.2,
                transform=ax.transAxes, clip_on=False)
            ax.add_patch(rect)
            fs = 5.5 if is_sec else 5.2
            fw = 'bold' if is_sec else 'normal'
            fc = DARK if is_sec else '#222222'
            ha = 'left' if ci == 0 else 'center'
            pad = 0.005 if ci == 0 else 0
            ax.text(x + pad, y + ROW_H/2, val[:80],
                    transform=ax.transAxes,
                    ha=ha, va='center', fontsize=fs, fontweight=fw,
                    color=fc, clip_on=True)

    plt.savefig(out_png, dpi=300, bbox_inches='tight', facecolor='white',
                edgecolor='none', pad_inches=0.02)
    plt.close()
    print(f"  Rendered table: {out_png}")


# ── Main ───────────────────────────────────────────────────────────────────────
print("Rendering table images...")
table1_png = rf"{OUT_DIR}\table1_rendered.png"
table2_png = rf"{OUT_DIR}\table2_rendered.png"

render_table1_png(rf"{OUT_DIR}\table1.xlsx", table1_png)
render_table2_png(rf"{OUT_DIR}\table2.xlsx", table2_png)

print("\nUpdating PPTX...")
prs = Presentation(PPTX_IN)
slides = prs.slides

# Slide 3 (index 2): Flowchart
replace_shape_with_image(slides[2], 'Picture 6',
                         rf"{OUT_DIR}\flowchart.png")

# Slide 4 (index 3): Table 1 (OLE → rendered image)
replace_shape_with_image(slides[3], 'Object 1', table1_png)

# Slide 5 (index 4): Table 2 (OLE → rendered image)
replace_shape_with_image(slides[4], 'Object 6', table2_png)

# Slide 6 (index 5): Figure 1
replace_shape_with_image(slides[5], 'Picture 3',
                         rf"{OUT_DIR}\fig1_hospice_by_year.png")

# Slide 7 (index 6): Figure 2
replace_shape_with_image(slides[6], 'Picture 3',
                         rf"{OUT_DIR}\fig2_days_io_to_death.png")

# Slide 8 (index 7): Figure 3
replace_shape_with_image(slides[7], 'Picture 3',
                         rf"{OUT_DIR}\fig3_hospice_los.png")

prs.save(PPTX_OUT)
print(f"\nSaved: {PPTX_OUT}")
